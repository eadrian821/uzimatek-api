"""
Corpus Ingestion Service
Indexes PDFs, PPTXs, and Markdown files into Qdrant for RAG
Medical corpus path: C:\\jarvis\\medical_corpus\\
Also handles Obsidian vault ingestion and OneNote local files
"""

import asyncio
import logging
import re
from pathlib import Path
from typing import List, Dict, Optional, Tuple
from datetime import datetime

import aiofiles

from app.core.config import settings
from app.services.memory import memory, embedder

logger = logging.getLogger(__name__)

# Corpus root
CORPUS = Path(settings.medical_corpus_path)

# Source books → specialty mapping (detected from filesystem scan)
SOURCE_MAP = {
    "kanski": "ophthalmology",
    "ophthalmol": "ophthalmology",
    "atlas of clinical ophthalmol": "ophthalmology",
    "review of ophthalmol": "ophthalmology",
    "dhingra": "ent",
    "ear nose throat": "ent",
    "ent questions": "ent",
    "gray": "anatomy",
    "grays anatomy": "anatomy",
    "first aid": "medicine",
    "guyton": "physiology",
    "hutchison": "medicine",
    "davidson": "medicine",
    "miller": "anaesthesia",
    "morgan": "anaesthesia",
    "stoelt": "anaesthesia",
    "lumb": "anaesthesia",
    "anaesthesia": "anaesthesia",
    "msc 500": "anaesthesia",
    "apley": "orthopaedics",
    "squires": "radiology",
    "dohnert": "radiology",
    "textbook of radiology": "radiology",
    "fitzpatrick": "dermatology",
    "habif": "dermatology",
    "derma": "dermatology",
    "junqueira": "physiology",
    "fundamentals": "medicine",
    "cawson": "dental",
    "oxford handbook of clinical dentistry": "dental",
    "oral surgery": "dental",
    "oral pathology": "dental",
    "atlas of diseases of the oral": "dental",
}


def classify_source(path: Path) -> str:
    """Classify a file into a medical specialty."""
    name = path.name.lower()
    parent = str(path.parent).lower()
    combined = f"{name} {parent}"
    for key, specialty in SOURCE_MAP.items():
        if key in combined:
            return specialty
    # Fallback: check parent folder name
    for part in path.parts:
        part_lower = part.lower()
        for specialty in ["ophthalmology", "ent", "anatomy", "medicine", "anaesthesia",
                          "orthopaedics", "radiology", "dermatology", "dental", "physiology"]:
            if specialty in part_lower:
                return specialty
    return "general"


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 100) -> List[str]:
    """
    Split text into overlapping chunks for optimal RAG retrieval.
    Medical facts often span sentences, so we use paragraph-aware splitting.
    """
    # Clean excessive whitespace
    text = re.sub(r'\n{3,}', '\n\n', text.strip())
    paragraphs = text.split('\n\n')
    chunks = []
    current = ""
    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        if len(current) + len(para) < chunk_size:
            current += "\n\n" + para if current else para
        else:
            if current:
                chunks.append(current)
                # Keep overlap from end of current chunk
                overlap_text = current[-overlap:] if len(current) > overlap else current
                current = overlap_text + "\n\n" + para
            else:
                # Single paragraph too long — split by sentence
                sentences = re.split(r'(?<=[.!?])\s+', para)
                for sent in sentences:
                    if len(current) + len(sent) < chunk_size:
                        current += " " + sent if current else sent
                    else:
                        if current:
                            chunks.append(current)
                        current = sent
    if current:
        chunks.append(current)
    # Filter out very short chunks
    return [c for c in chunks if len(c) > 50]


class CorpusIngestionService:
    """Ingests all medical source materials into Qdrant."""

    # ─── PDF Ingestion ────────────────────────────────────────────────────────

    async def ingest_pdf(self, pdf_path: Path, force: bool = False) -> int:
        """Extract, chunk, and embed a PDF into medical_gold_standard."""
        marker_path = CORPUS / "_indexed" / f"{pdf_path.stem}.indexed"
        if marker_path.exists() and not force:
            logger.debug(f"Skipping already indexed: {pdf_path.name}")
            return 0
        logger.info(f"Ingesting PDF: {pdf_path.name}")
        try:
            import fitz  # PyMuPDF is much better at extracting from medical/scanned PDFs
            doc = fitz.open(str(pdf_path))
            full_text = ""
            for page in doc:
                try:
                    text = page.get_text()
                    if text:
                        full_text += text + "\n\n"
                except Exception:
                    continue
            doc.close()
            
            # OCR Fallback for scanned/image-only PDFs
            if not full_text.strip():
                logger.info(f"No text found in {pdf_path.name}, attempting OCR fallback...")
                try:
                    import pytesseract
                    from pdf2image import convert_from_path
                    
                    # Convert PDF pages to images (using poppler)
                    images = convert_from_path(str(pdf_path), dpi=300)
                    for i, img in enumerate(images):
                        logger.info(f"Running OCR on {pdf_path.name} page {i+1}...")
                        text = pytesseract.image_to_string(img)
                        if text:
                            full_text += text + "\n\n"
                except Exception as e:
                    logger.error(f"OCR fallback failed for {pdf_path.name}: {e}")
                    
            if not full_text.strip():
                logger.warning(f"No text extracted from {pdf_path.name} even after OCR attempt.")
                return 0
            specialty = classify_source(pdf_path)
            chunks = chunk_text(full_text)
            metadata = {
                "source_file": pdf_path.name,
                "specialty": specialty,
                "source_type": "textbook",
                "ingested_at": datetime.now().isoformat()
            }
            count = await memory.semantic.store_chunked_document(
                collection="medical_gold_standard",
                source_path=str(pdf_path),
                chunks=chunks,
                base_metadata=metadata
            )
            # Mark as indexed
            marker_path.parent.mkdir(parents=True, exist_ok=True)
            marker_path.write_text(f"indexed {count} chunks at {datetime.now().isoformat()}")
            logger.info(f"✅ Indexed {count} chunks from {pdf_path.name} [{specialty}]")
            return count
        except Exception as e:
            logger.error(f"PDF ingestion error for {pdf_path.name}: {e}")
            return 0

    # ─── PPTX Ingestion ───────────────────────────────────────────────────────

    async def ingest_pptx(self, pptx_path: Path, force: bool = False) -> int:
        """Extract text from rotation slides and index to medical_gold_standard."""
        marker_path = CORPUS / "_indexed" / f"{pptx_path.stem}.indexed"
        if marker_path.exists() and not force:
            return 0
        logger.info(f"Ingesting PPTX: {pptx_path.name}")
        try:
            from pptx import Presentation
            from pptx.util import Inches
            prs = Presentation(str(pptx_path))
            slides_text = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    slides_text.append(f"[Slide {i+1}] " + " | ".join(slide_text))
            if not slides_text:
                return 0
            full_text = "\n\n".join(slides_text)
            specialty = classify_source(pptx_path)
            chunks = chunk_text(full_text, chunk_size=400)  # Smaller for slides
            metadata = {
                "source_file": pptx_path.name,
                "specialty": specialty,
                "source_type": "lecture_slide",
                "ingested_at": datetime.now().isoformat()
            }
            count = await memory.semantic.store_chunked_document(
                collection="medical_gold_standard",
                source_path=str(pptx_path),
                chunks=chunks,
                base_metadata=metadata
            )
            marker_path.parent.mkdir(parents=True, exist_ok=True)
            marker_path.write_text(f"indexed {count} chunks at {datetime.now().isoformat()}")
            logger.info(f"✅ Indexed {count} slide chunks from {pptx_path.name} [{specialty}]")
            return count
        except Exception as e:
            logger.error(f"PPTX ingestion error: {e}")
            return 0

    # ─── Markdown (Obsidian vault) Ingestion ─────────────────────────────────

    async def ingest_markdown(self, md_path: Path) -> int:
        """Index a vault markdown note to obsidian_notes collection."""
        try:
            async with aiofiles.open(md_path, 'r', encoding='utf-8') as f:
                content = await f.read()
            if len(content) < 100:
                return 0
            # Strip Obsidian-specific syntax for cleaner embedding
            clean = re.sub(r'!\[.*?\]\(.*?\)', '', content)  # Remove image links
            clean = re.sub(r'\[\[.*?\]\]', lambda m: m.group(0)[2:-2], clean)  # Unwrap wiki-links
            clean = re.sub(r'---\n.*?---\n', '', clean, flags=re.DOTALL)  # Remove frontmatter
            chunks = chunk_text(clean, chunk_size=600)
            metadata = {
                "source_file": md_path.name,
                "source_type": "obsidian_note",
                "vault_path": str(md_path),
                "modified_at": datetime.fromtimestamp(md_path.stat().st_mtime).isoformat()
            }
            count = await memory.semantic.store_chunked_document(
                collection="obsidian_notes",
                source_path=str(md_path),
                chunks=chunks,
                base_metadata=metadata
            )
            return count
        except Exception as e:
            logger.error(f"Markdown ingestion error for {md_path.name}: {e}")
            return 0

    # ─── OneNote Local File Ingestion ─────────────────────────────────────────

    async def ingest_onenote_folder(self, onenote_path: Path = None) -> int:
        """Parse local OneNote .one files and index to onenote_annotations."""
        base = onenote_path or Path(settings.onenote_path)
        if not base.exists():
            logger.warning(f"OneNote path not found: {base}")
            return 0
        total = 0
        # OneNote local files are .one format (proprietary) — extract via file text
        for one_file in base.rglob("*.one"):
            try:
                total += await self._parse_one_file(one_file)
            except Exception as e:
                logger.warning(f"OneNote parse error {one_file.name}: {e}")
        # Also check for exported .docx/.txt from OneNote
        for exported in base.rglob("*.docx"):
            try:
                from docx import Document
                doc = Document(str(exported))
                text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                if len(text) > 100:
                    chunks = chunk_text(text)
                    await memory.semantic.store_chunked_document(
                        collection="onenote_annotations",
                        source_path=str(exported),
                        chunks=chunks,
                        base_metadata={"source_type": "onenote_export", "source_file": exported.name}
                    )
                    total += len(chunks)
            except Exception:
                continue
        logger.info(f"OneNote ingestion complete: {total} chunks")
        return total

    async def _parse_one_file(self, path: Path) -> int:
        """Extract readable text from .one binary format."""
        try:
            import olefile
            if not olefile.isOleFile(str(path)):
                return 0
            ole = olefile.OleFileIO(str(path))
            text_parts = []
            for entry in ole.listdir():
                try:
                    stream = ole.openstream(entry)
                    raw = stream.read()
                    # Extract printable ASCII/UTF-8 sequences
                    decoded = raw.decode('utf-8', errors='ignore')
                    printable = re.sub(r'[^\x20-\x7E\n\t]', ' ', decoded)
                    cleaned = re.sub(r' {3,}', ' ', printable).strip()
                    if len(cleaned) > 50:
                        text_parts.append(cleaned)
                except Exception:
                    continue
            ole.close()
            if not text_parts:
                return 0
            full_text = "\n\n".join(text_parts[:20])  # Limit to top sections
            chunks = chunk_text(full_text)
            await memory.semantic.store_chunked_document(
                collection="onenote_annotations",
                source_path=str(path),
                chunks=chunks,
                base_metadata={"source_type": "onenote_local", "source_file": path.name}
            )
            # Check for #card tags → flag for atomization
            card_lines = [c for c in full_text.split('\n') if '#card' in c.lower()]
            for line in card_lines:
                await memory.log_recommendation(line, source="onenote_card_tag")
            return len(chunks)
        except Exception as e:
            logger.warning(f"OLE parse error: {e}")
            return 0

    # ─── Full Corpus Bootstrap ────────────────────────────────────────────────

    async def bootstrap_corpus(self, source_root: Path = None) -> Dict[str, int]:
        """
        Index ALL medical materials found on the filesystem.
        Searches across known paths from system scan.
        """
        search_paths = [
            Path(r"C:\Users\eadri\Workspace\MedSchool"),
            Path(r"C:\Users\eadri\Desktop\Dental"),
            Path(r"C:\Users\eadri\Downloads"),
            Path(r"C:\Users\eadri\Desktop\Derma"),
            CORPUS
        ]
        if source_root:
            search_paths.insert(0, source_root)

        stats = {"pdfs": 0, "pptx": 0, "markdown": 0, "total_chunks": 0}

        for search_path in search_paths:
            if not search_path.exists():
                continue
            logger.info(f"Scanning: {search_path}")

            # PDFs
            for pdf in search_path.rglob("*.pdf"):
                # Skip system/temp files
                if any(skip in str(pdf).lower() for skip in ["venv", "node_modules", ".git", "temp"]):
                    continue
                if pdf.stat().st_size > 100:  # Skip empty files
                    count = await self.ingest_pdf(pdf)
                    if count:
                        stats["pdfs"] += 1
                        stats["total_chunks"] += count
                    await asyncio.sleep(0.1)  # Rate limit embedding API

            # PPTXs
            for pptx in search_path.rglob("*.pptx"):
                # Skip temporary PowerPoint lock files
                if pptx.name.startswith("~$"):
                    continue
                if pptx.stat().st_size > 100:
                    count = await self.ingest_pptx(pptx)
                    if count:
                        stats["pptx"] += 1
                        stats["total_chunks"] += count
                    await asyncio.sleep(0.1)

        # Obsidian vault
        vault = Path(settings.obsidian_vault_path)
        for md in vault.rglob("*.md"):
            if ".obsidian" not in str(md):
                count = await self.ingest_markdown(md)
                stats["markdown"] += 1
                stats["total_chunks"] += count
                await asyncio.sleep(0.05)

        logger.info(f"Corpus bootstrap complete: {stats}")
        return stats

    async def ingest_single_file(self, path: Path) -> int:
        """Ingest a single file — called when vault watcher detects changes."""
        path = Path(path)
        if path.suffix.lower() == ".pdf":
            return await self.ingest_pdf(path, force=True)
        elif path.suffix.lower() == ".pptx":
            return await self.ingest_pptx(path, force=True)
        elif path.suffix.lower() == ".md":
            return await self.ingest_markdown(path)
        return 0


# Global instance
corpus = CorpusIngestionService()
