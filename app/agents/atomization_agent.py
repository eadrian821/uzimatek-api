"""
Atomization Agent
Converts Delta Engine gaps into Anki cards with clinical images (dual-coding)
Also handles: PPTX slides → cards, OneNote #card tags → cards
"""

import logging
import re
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from app.agents.base import BaseAgent, AgentResponse
from app.core.config import settings
from app.services.anki_service import anki, image_fetcher
from app.services.memory import memory

logger = logging.getLogger(__name__)


class AtomizationAgent(BaseAgent):
    """
    Atomization pipeline: Gap text → Anki card(s) + clinical image.
    
    Design: One gap → one cloze + one basic card pair.
    Images added for anatomy, imaging, and pathology topics.
    Cards tagged with rotation, source, and date for FSRS tracking.
    """

    # Topics where visual dual-coding is most valuable
    VISUAL_TOPICS = [
        "anatomy", "imaging", "x-ray", "mri", "ct", "fundus", "skin", "rash",
        "ear", "eye", "nose", "throat", "cardiac", "ecg", "fracture", "pathology",
        "histology", "microscopy", "tympanic", "retina", "glaucoma", "cholesteatoma",
        "dental", "oral", "bone", "joint", "wound", "surgical", "lesion"
    ]

    def __init__(self):
        super().__init__(name="atomization", description="Gap to Anki card atomization")

    @property
    def system_prompt(self) -> str:
        return f"""You are JARVIS's Atomization Agent for {settings.user_name}, a medical student.

Your function: Convert a clinical knowledge gap into high-quality Anki flashcards.

Rules for card creation:
1. ONE atomic fact per card — never put multiple facts on one card
2. For mechanisms: use cloze deletion ("The mechanism of X is {{{{c1::Y}}}}")
3. For clinical correlations: use basic Q&A format
4. Front must be a genuine question the student can test themselves against
5. Back must be clinically precise — include drug dosages, staging criteria, specific values where relevant
6. Add a memory mnemonic or clinical pearl in "Extra" if one exists
7. Never use vague language — "some of" or "often" are forbidden

Output JSON format:
{{
    "cards": [
        {{
            "type": "cloze",
            "content": "The first-line treatment for X is {{{{c1::drug_name}}}} at dose {{{{c2::dose}}}}",
            "extra": "Pearl: ...",
            "needs_image": true/false,
            "image_topic": "topic keyword for image search"
        }},
        {{
            "type": "basic",
            "front": "What is the pathophysiology of X?",
            "back": "Precise mechanism...",
            "extra": "Mnemonic: ...",
            "needs_image": false,
            "image_topic": ""
        }}
    ]
}}
"""

    @property
    def capabilities(self) -> List[str]:
        return ["gap_atomization", "slide_atomization", "onenote_atomization"]

    async def process(
        self,
        message: str,
        intent: Any,
        context: Any,
        attachments: List[Dict] = None
    ) -> AgentResponse:
        # Check if this is a slide file
        if attachments:
            for att in attachments:
                if att.get("type") == "document" and att.get("file_name", "").endswith(".pptx"):
                    return await self._handle_slide_file(att, context)
        return await self.atomize_gaps_from_text(message, context=context)

    async def atomize_gaps(
        self,
        gaps: List[str],
        topic: str,
        rotation: str = "General",
        deck: str = None
    ) -> AgentResponse:
        """Main entry point: convert a list of gap strings into Anki cards."""
        if not gaps:
            return AgentResponse(agent=self.name, content="No gaps provided to atomize.", confidence=0.5)

        # Check AnkiConnect availability
        if not await anki.is_available():
            return AgentResponse(
                agent=self.name,
                content=(
                    "⚠️ **AnkiConnect not running.**\n\n"
                    "1. Open Anki on your computer\n"
                    "2. Go to Tools → Add-ons → Browse & Install\n"
                    "3. Enter code: `2055492159`\n"
                    "4. Restart Anki\n\n"
                    "Then try again with `/anki`"
                ),
                confidence=0.9
            )

        target_deck = deck or f"{settings.anki_default_deck}::{rotation}::{topic.replace(' ', '_')}"
        tags = [
            f"rotation::{rotation.lower().replace(' ', '_')}",
            f"topic::{topic.lower().replace(' ', '_')}",
            "source::delta_engine",
            f"date::{datetime.now().strftime('%Y-%m')}",
            "jarvis"
        ]

        all_notes = []
        for gap in gaps:
            cards = await self._generate_cards_for_gap(gap, topic)
            for card in cards:
                card["deck"] = target_deck
                card["tags"] = tags
                card["source"] = f"{topic} — Delta Engine"
                all_notes.append(card)

        if not all_notes:
            return AgentResponse(agent=self.name, content="Card generation failed.", confidence=0.2)

        # Fetch images for visual topics
        all_notes = await self._add_images(all_notes, topic)

        # Push to Anki
        success, fail = await anki.add_batch(all_notes)
        await anki.sync()

        msg = (
            f"🃏 **Cards Pushed to Anki**\n\n"
            f"✅ {success} cards added to `{target_deck}`\n"
            f"{'❌ ' + str(fail) + ' failed' if fail else ''}\n\n"
            f"Tags: {', '.join(tags[:3])}\n"
            f"_FSRS will schedule your first review automatically._"
        )
        return AgentResponse(
            agent=self.name,
            content=msg,
            confidence=0.95,
            actions_taken=[{"action": "anki_push", "cards": success, "deck": target_deck}],
            data={"success": success, "fail": fail, "deck": target_deck}
        )

    async def atomize_gaps_from_text(self, text: str, context: Any = None) -> AgentResponse:
        """Parse Gap:: lines from text and atomize them."""
        gaps = []
        topic = "General"
        for line in text.split('\n'):
            line = line.strip()
            if line.startswith("Gap ::"):
                gaps.append(line.replace("Gap ::", "").strip())
            elif "topic:" in line.lower():
                topic = line.split(":", 1)[1].strip()
        if not gaps:
            # Treat entire text as a single concept to atomize
            gaps = [text]
        return await self.atomize_gaps(gaps, topic=topic)

    async def atomize_slide_deck(self, pptx_path: Path, rotation: str = "General") -> AgentResponse:
        """Process a PPTX file and generate Anki cards from each slide."""
        try:
            from pptx import Presentation
            prs = Presentation(str(pptx_path))
        except Exception as e:
            return AgentResponse(agent=self.name, content=f"Could not read PPTX: {e}", confidence=0.2)

        topic = pptx_path.stem.replace('_', ' ')
        all_notes = []
        tags = [
            f"rotation::{rotation.lower()}",
            f"source::slide_{pptx_path.stem.lower()[:20]}",
            f"date::{datetime.now().strftime('%Y-%m')}",
            "jarvis"
        ]
        target_deck = f"{settings.anki_default_deck}::{rotation}::Slides"

        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if len(slide_text) < 2:
                continue  # Skip near-empty slides
            slide_content = " | ".join(slide_text)
            if len(slide_content) < 30:
                continue
            # Generate cards for this slide
            slide_gaps = [f"From slide {i+1}: {slide_content[:300]}"]
            cards = await self._generate_cards_for_gap(slide_gaps[0], topic)
            for card in cards[:2]:  # Max 2 cards per slide
                card["deck"] = target_deck
                card["tags"] = tags
                card["source"] = f"{pptx_path.name} — Slide {i+1}"
                all_notes.append(card)

        all_notes = await self._add_images(all_notes, topic)
        if not await anki.is_available():
            return AgentResponse(agent=self.name, content="⚠️ AnkiConnect not running. Open Anki first.", confidence=0.9)

        success, fail = await anki.add_batch(all_notes)
        await anki.sync()

        return AgentResponse(
            agent=self.name,
            content=f"🃏 **Slide Deck Atomized**\n\n✅ {success} cards from `{pptx_path.name}` → `{target_deck}`\n{'❌ ' + str(fail) + ' failed' if fail else ''}",
            confidence=0.9,
            actions_taken=[{"action": "slide_atomization", "file": pptx_path.name, "cards": success}]
        )

    # ─── Internal Helpers ─────────────────────────────────────────────────────

    async def _generate_cards_for_gap(self, gap_text: str, topic: str) -> List[Dict]:
        """Use Claude to generate cloze + basic card pair for a gap."""
        prompt = f"""Generate Anki cards for this clinical gap:

TOPIC: {topic}
GAP: {gap_text}

Return JSON only. Create 1 cloze + 1 basic card."""

        try:
            response = await self._call_claude(
                messages=[{"role": "user", "content": prompt}],
                max_tokens=800
            )
            # Extract JSON
            json_match = re.search(r'\{.*\}', response, re.DOTALL)
            if json_match:
                import json
                data = json.loads(json_match.group())
                cards = data.get("cards", [])
                # Convert to add_batch format
                result = []
                for card in cards:
                    if card.get("type") == "cloze":
                        result.append({
                            "type": "cloze",
                            "content": card.get("content", gap_text),
                            "extra": card.get("extra", ""),
                            "needs_image": card.get("needs_image", False),
                            "image_topic": card.get("image_topic", topic)
                        })
                    else:
                        result.append({
                            "type": "basic",
                            "front": card.get("front", f"What is the key fact about {topic}?"),
                            "back": card.get("back", gap_text),
                            "extra": card.get("extra", ""),
                            "needs_image": card.get("needs_image", False),
                            "image_topic": card.get("image_topic", topic)
                        })
                return result
        except Exception as e:
            logger.warning(f"Card generation error: {e}")

        # Fallback: create simple basic card from gap
        return [{
            "type": "basic",
            "front": f"[{topic}] {gap_text[:100]}",
            "back": gap_text,
            "extra": "",
            "needs_image": any(kw in gap_text.lower() for kw in self.VISUAL_TOPICS),
            "image_topic": topic
        }]

    async def _add_images(self, notes: List[Dict], default_topic: str) -> List[Dict]:
        """Fetch and embed clinical images for cards that need them."""
        for note in notes:
            needs_img = note.get("needs_image", False)
            topic = note.get("image_topic", default_topic)
            # Auto-detect if image would be valuable
            content = note.get("content") or note.get("back") or note.get("front", "")
            if not needs_img:
                needs_img = any(kw in content.lower() for kw in self.VISUAL_TOPICS)
            if needs_img:
                img_html = await image_fetcher.get_image_for_card(topic, content)
                note["image_html"] = img_html or ""
            else:
                note["image_html"] = ""
        return notes

    async def _handle_slide_file(self, attachment: Dict, context: Any) -> AgentResponse:
        """Handle PPTX file sent via Telegram."""
        # This would normally download the file first
        return AgentResponse(
            agent=self.name,
            content="📊 Send the PPTX file and I'll atomize it into Anki cards with images. Use `/anki-slides` command.",
            confidence=0.8
        )
